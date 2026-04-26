"""
Generate EMBRACE AI — 12-Month Initiative PowerPoint Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "EMBRACE_AI_Presentation.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Brand Colours ───
PETROL = RGBColor(0x00, 0x99, 0x99)
PETROL_DARK = RGBColor(0x00, 0x7A, 0x7A)
DARK = RGBColor(0x00, 0x00, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF7, 0xF7, 0xF8)
GREY = RGBColor(0x6E, 0x6E, 0x73)
GREY_LIGHT = RGBColor(0xA1, 0xA1, 0xA6)
GREEN = RGBColor(0x00, 0xD4, 0xA0)
SPARK_COL = RGBColor(0xF5, 0x9E, 0x0B)
BUILD_COL = RGBColor(0x3B, 0x82, 0xF6)
APPLY_COL = RGBColor(0xF9, 0x73, 0x16)
DELIVER_COL = RGBColor(0xEF, 0x44, 0x44)

# ─── Helpers ───
def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri', spacing=1.1):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    p.line_spacing = spacing
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                       color=WHITE, bold=False, bullet=False, line_colors=None, line_bolds=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "•  " if bullet else ""
        p.text = prefix + line
        p.font.size = Pt(font_size)
        p.font.color.rgb = (line_colors[i] if line_colors and i < len(line_colors) else color)
        p.font.bold = (line_bolds[i] if line_bolds and i < len(line_bolds) else bold)
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.line_spacing = 1.2
    return txBox

def add_divider(slide, left, top, width, color=PETROL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def slide_footer(slide, text="EMBRACE AI  ·  Engineering Systems  ·  Siemens"):
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(8), Inches(0.4),
                 text, font_size=9, color=GREY_LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)

# Left accent bar
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PETROL)

# Siemens text logo
add_text_box(slide, Inches(1.0), Inches(0.8), Inches(4), Inches(0.6),
             "SIEMENS", font_size=22, color=PETROL, bold=True)

add_divider(slide, Inches(1.0), Inches(1.5), Inches(3))

add_text_box(slide, Inches(1.0), Inches(2.2), Inches(8), Inches(1.2),
             "EMBRACE AI", font_size=60, color=WHITE, bold=True)

add_text_box(slide, Inches(1.0), Inches(3.5), Inches(8), Inches(0.8),
             "12-Month AI Initiative Plan", font_size=28, color=GREEN)

add_text_box(slide, Inches(1.0), Inches(4.5), Inches(8), Inches(0.6),
             "Engineering Systems  ·  FT D AA IN SGI DET ENGSYS", font_size=16, color=GREY)

add_text_box(slide, Inches(1.0), Inches(5.2), Inches(8), Inches(0.5),
             "Siemens, Chennai, India", font_size=14, color=GREY_LIGHT)

# Right side — decorative circles
for i, (x, y, size, opacity_color) in enumerate([
    (10.0, 1.0, 2.5, RGBColor(0x00, 0x50, 0x50)),
    (9.5, 3.5, 1.8, RGBColor(0x00, 0x40, 0x40)),
    (11.0, 4.5, 1.2, RGBColor(0x00, 0x60, 0x60)),
]):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = opacity_color
    shape.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(6.2), Inches(6), Inches(0.8),
             "Prepared by: Vigneshvar SA\nvigneshvar.sa@siemens.com  ·  April 2025",
             font_size=12, color=GREY_LIGHT, spacing=1.4)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.3), DARK)
add_text_box(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.7),
             "Agenda", font_size=36, color=WHITE, bold=True)

agenda_items = [
    "About Engineering Systems",
    "Initiative Overview & Objectives",
    "The Team — 16 Members, 4 Domains",
    "Syllabus Design — 4-Phase Maturity Model",
    "Points & Rewards System",
    "12-Month Detailed Plan (Month by Month)",
    "The Dashboard — Live Demo",
    "Expected Outcomes & Deliverables",
]
for i, item in enumerate(agenda_items):
    y = 1.8 + i * 0.6
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y), Inches(0.45), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PETROL
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = 'Calibri'

    add_text_box(slide, Inches(1.7), Inches(y + 0.03), Inches(8), Inches(0.45),
                 item, font_size=20, color=DARK)

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — ABOUT ENGINEERING SYSTEMS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.3), DARK)
add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.7),
             "About Engineering Systems", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.6),
             "The digital and operational backbone for our North American partners — managing critical PLM, Manufacturing, and OT applications.",
             font_size=15, color=GREY)

# Three pillar cards
pillars = [
    ("PLM", "Product Lifecycle Management", PETROL,
     ["Teamcenter Admin & Customization", "MCAD (Creo, NX Open)", "ECAD (Zuken e3)",
      "SAP Integration (T4S)", "Rulestream Support"]),
    ("MFG", "Manufacturing", SPARK_COL,
     ["OpCenter ED (MES)", "OpCenter OEE", "Praxis Management",
      "AGW Administration", "DMM Administration"]),
    ("OT", "OT & Automation", RGBColor(0x10, 0xB9, 0x81),
     ["SINEC NMS Management", "Data Analytics", "SCADA & PLC Systems",
      "Industrial Edge & IIH", "AI / GenAI (this initiative!)"]),
]

for i, (label, title, color, items) in enumerate(pillars):
    x = 0.8 + i * 4.1
    card = add_shape(slide, Inches(x), Inches(2.6), Inches(3.8), Inches(4.2), LIGHT)
    # Color top bar
    add_shape(slide, Inches(x), Inches(2.6), Inches(3.8), Inches(0.12), color)
    add_text_box(slide, Inches(x + 0.3), Inches(2.9), Inches(3.2), Inches(0.5),
                 label, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(x + 0.3), Inches(3.25), Inches(3.2), Inches(0.4),
                 title, font_size=12, color=GREY)
    add_multiline_box(slide, Inches(x + 0.3), Inches(3.7), Inches(3.2), Inches(3.0),
                       items, font_size=13, color=DARK, bullet=True)

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — INITIATIVE OVERVIEW
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Initiative Overview", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.3), Inches(3))

# Key stats in cards
stats = [
    ("16", "Team Members", "Across 4 domains"),
    ("12", "Monthly Sessions", "April 2025 – March 2026"),
    ("4", "Phases", "SPARK → BUILD → APPLY → DELIVER"),
    ("🏆", "Star Points", "Top 3 win gift cards & vouchers"),
]
for i, (num, label, desc) in enumerate(stats):
    x = 0.8 + i * 3.1
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(2.8), Inches(1.8), RGBColor(0x00, 0x20, 0x40))
    add_text_box(slide, Inches(x + 0.3), Inches(1.95), Inches(2.2), Inches(0.8),
                 str(num), font_size=40, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.3), Inches(2.65), Inches(2.2), Inches(0.35),
                 label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.3), Inches(3.0), Inches(2.2), Inches(0.45),
                 desc, font_size=11, color=GREY, alignment=PP_ALIGN.CENTER)

# Objectives
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.5),
             "OBJECTIVES", font_size=14, color=PETROL, bold=True)

objectives = [
    "Build GenAI literacy across the entire team — from curiosity to competence",
    "Transition from fun learning to real business-value delivery over 12 months",
    "Produce 16 documented AI use-cases with working prototypes by year-end",
    "Create a replicable playbook for other Siemens teams in India",
    "Measurable productivity gains — aggregate time saved, quality improved",
]
add_multiline_box(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(2.5),
                   objectives, font_size=15, color=WHITE, bullet=True)

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — THE TEAM
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.3), DARK)
add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.7),
             "The Team — 16 Members", font_size=36, color=WHITE, bold=True)

members = [
    ("Aravinthan Dhanabal", "MFG"), ("Aswin Gangadharan", "OT"),
    ("Balaji M", "PLM"), ("Dharankumar N", "PLM"),
    ("Gopalakrishnan R", "PLM"), ("John Yabesh Johnson", "MFG"),
    ("Minal R Nilange", "MFG"), ("Pandiyarajan N", "Cross"),
    ("Priyadharshani VS", "PLM"), ("Rizwan Ali M", "PLM"),
    ("Saravanan J", "Cross"), ("Sivakumar D", "Cross"),
    ("Soniya Dhayalan", "OT"), ("Tippu Sultan Shaik", "PLM"),
    ("Vidya Srivatsan", "Cross"), ("Vigneshvar SA ★", "PLM"),
]

domain_colors = {"PLM": PETROL, "MFG": SPARK_COL, "OT": RGBColor(0x10, 0xB9, 0x81), "Cross": RGBColor(0x8B, 0x5C, 0xF6)}
domain_labels = {"PLM": "PLM", "MFG": "Manufacturing", "OT": "OT & Automation", "Cross": "Cross-Functional"}

for i, (name, domain) in enumerate(members):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.1
    y = 1.7 + row * 1.35
    card = add_shape(slide, Inches(x), Inches(y), Inches(2.8), Inches(1.15), LIGHT)
    # Avatar circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(y + 0.2), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = domain_colors[domain]
    circle.line.fill.background()
    initials = ''.join(w[0] for w in name.replace('★','').strip().split()[:2])
    tf = circle.text_frame
    tf.paragraphs[0].text = initials
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = 'Calibri'

    add_text_box(slide, Inches(x + 0.8), Inches(y + 0.15), Inches(1.8), Inches(0.35),
                 name, font_size=12, color=DARK, bold=True)
    add_text_box(slide, Inches(x + 0.8), Inches(y + 0.5), Inches(1.8), Inches(0.3),
                 domain_labels[domain], font_size=10, color=domain_colors[domain])

# Domain legend
for i, (key, label) in enumerate(domain_labels.items()):
    x = 0.8 + i * 3.1
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(7.0), Inches(0.2), Inches(0.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = domain_colors[key]
    circle.line.fill.background()
    add_text_box(slide, Inches(x + 0.3), Inches(6.95), Inches(2), Inches(0.3),
                 label, font_size=10, color=GREY)

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — 4-PHASE MATURITY MODEL
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.3), DARK)
add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.7),
             "Syllabus Design — 4-Phase Maturity Model", font_size=36, color=WHITE, bold=True)

phases = [
    ("SPARK", "Months 1–3", "Curiosity & Fun", "80% Fun / 20% Learning",
     SPARK_COL, ["AI Portraits", "AI Storyteller", "Prompt Battle"]),
    ("BUILD", "Months 4–6", "Skill Building", "50% Fun / 50% Learning",
     BUILD_COL, ["Copilot Power User", "Code with AI", "Smart Docs / RAG"]),
    ("APPLY", "Months 7–9", "Domain Automation", "20% Fun / 80% Business",
     APPLY_COL, ["Workflow Automation", "Agents Hackathon", "Data Detective"]),
    ("DELIVER", "Months 10–12", "Business Value", "100% Business Impact",
     DELIVER_COL, ["Use Case: Define", "Use Case: Build", "Grand Finale"]),
]

for i, (name, months, theme, orientation, color, activities) in enumerate(phases):
    x = 0.5 + i * 3.2
    # Card
    card = add_shape(slide, Inches(x), Inches(1.7), Inches(2.9), Inches(5.2), LIGHT)
    # Color top bar
    add_shape(slide, Inches(x), Inches(1.7), Inches(2.9), Inches(0.15), color)
    # Phase name
    add_text_box(slide, Inches(x + 0.25), Inches(2.0), Inches(2.4), Inches(0.5),
                 name, font_size=28, color=color, bold=True)
    add_text_box(slide, Inches(x + 0.25), Inches(2.5), Inches(2.4), Inches(0.3),
                 months, font_size=13, color=GREY, bold=True)
    add_text_box(slide, Inches(x + 0.25), Inches(2.85), Inches(2.4), Inches(0.3),
                 theme, font_size=14, color=DARK, bold=True)
    add_text_box(slide, Inches(x + 0.25), Inches(3.2), Inches(2.4), Inches(0.3),
                 orientation, font_size=11, color=GREY)
    # Divider
    add_shape(slide, Inches(x + 0.25), Inches(3.6), Inches(2.4), Pt(1), RGBColor(0xE5, 0xE5, 0xEA))
    # Activities
    add_multiline_box(slide, Inches(x + 0.25), Inches(3.8), Inches(2.4), Inches(2.5),
                       activities, font_size=13, color=DARK, bullet=True)

    # Arrow between phases
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 3.0), Inches(3.8), Inches(0.35), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GREY_LIGHT
        arrow.line.fill.background()

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — POINTS & REWARDS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Points & Rewards System", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.3), Inches(3))

# Points table
points_data = [
    ("Seminar Presenter (peer-rated)", "Up to 5 pts"),
    ("Quiz Winner — 1st / 2nd / 3rd", "3 / 2 / 1 pts"),
    ("Activity Completion (on time)", "3 pts"),
    ("Activity Category Awards", "5 pts each"),
    ("First to Submit Bonus", "2 pts"),
    ("People's Choice (peer vote)", "3 pts"),
    ("Grand Finale — Best Use Case", "15 pts"),
    ("Grand Finale — Runner-Up", "10 pts"),
    ("Grand Finale — 2nd Runner-Up", "7 pts"),
]

# Header
add_shape(slide, Inches(0.8), Inches(1.7), Inches(7.5), Inches(0.55), PETROL)
add_text_box(slide, Inches(1.0), Inches(1.75), Inches(5), Inches(0.45),
             "Category", font_size=14, color=WHITE, bold=True)
add_text_box(slide, Inches(6.0), Inches(1.75), Inches(2), Inches(0.45),
             "Points", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

for i, (cat, pts) in enumerate(points_data):
    y = 2.3 + i * 0.48
    bg_color = RGBColor(0x00, 0x20, 0x40) if i % 2 == 0 else RGBColor(0x00, 0x18, 0x38)
    add_shape(slide, Inches(0.8), Inches(y), Inches(7.5), Inches(0.45), bg_color)
    add_text_box(slide, Inches(1.0), Inches(y + 0.05), Inches(5), Inches(0.35),
                 cat, font_size=13, color=WHITE)
    add_text_box(slide, Inches(6.0), Inches(y + 0.05), Inches(2), Inches(0.35),
                 pts, font_size=13, color=GREEN, bold=True, alignment=PP_ALIGN.RIGHT)

# Right side — Prize card
card = add_shape(slide, Inches(9.0), Inches(1.7), Inches(3.8), Inches(3.5), RGBColor(0x00, 0x20, 0x40))
add_text_box(slide, Inches(9.3), Inches(1.9), Inches(3.2), Inches(0.5),
             "🏆  YEAR-END PRIZES", font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(9.3), Inches(2.5), Inches(3.2), Pt(1), RGBColor(0x00, 0x40, 0x40))

prizes = [
    ("🥇  1st Place", "Star Points → Gift Cards"),
    ("🥈  2nd Place", "Star Points → Vouchers"),
    ("🥉  3rd Place", "Star Points → Coupons"),
    ("📜  All 16", "Certificate of Completion"),
]
for i, (label, desc) in enumerate(prizes):
    y = 2.7 + i * 0.6
    add_text_box(slide, Inches(9.3), Inches(y), Inches(3.2), Inches(0.3),
                 label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.3), Inches(y + 0.3), Inches(3.2), Inches(0.25),
                 desc, font_size=11, color=GREY, alignment=PP_ALIGN.CENTER)

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — MONTHLY STRUCTURE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.3), DARK)
add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.7),
             "Recurring Monthly Structure", font_size=36, color=WHITE, bold=True)

segments = [
    ("🎤", "Seminar", "15–20 min", "1 person presents a drawn table topic. Prepared over previous month. Followed by 5-question quiz.", PETROL),
    ("🎯", "Main Activity", "45–90 min", "Either spontaneous (in-session) or pre-worked (submitted before). Core skill-building exercise.", BUILD_COL),
    ("🏆", "Awards & Leaderboard", "10 min", "Points tallied. Leaderboard updated live on the dashboard. Category winners announced. Next topic drawn.", SPARK_COL),
]

for i, (icon, title, duration, desc, color) in enumerate(segments):
    y = 1.8 + i * 1.7
    # Number/icon area
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.1), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(2.2), Inches(y), Inches(3), Inches(0.4),
                 title, font_size=22, color=DARK, bold=True)
    add_text_box(slide, Inches(2.2), Inches(y + 0.4), Inches(1.5), Inches(0.3),
                 duration, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(2.2), Inches(y + 0.7), Inches(8), Inches(0.5),
                 desc, font_size=14, color=GREY)

    # Connector line
    if i < 2:
        add_shape(slide, Inches(1.38), Inches(y + 1.0), Pt(2), Inches(0.7), RGBColor(0xE5, 0xE5, 0xEA))

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDES 9–20 — INDIVIDUAL MONTH DETAILS
# ═══════════════════════════════════════════════════════════
months_detail = [
    {
        "num": 1, "title": "Who's Who: AI Edition", "phase": "SPARK", "color": SPARK_COL,
        "date": "Apr 8, 2025", "type": "Pre-worked", "dur": "90 min",
        "desc": "Each member crafts AI prompts and generates portraits + one-liners for every colleague. Fun kickoff activity.",
        "seminar": "What is GenAI and Why Should Engineers Care?",
        "presenter": "Vigneshvar SA",
        "awards": ["Most Creative Prompt", "Best AI Image", "Funniest One-Liner", "First to Complete"],
        "skills": ["Prompt engineering basics", "AI image generation", "Iterative prompting"],
    },
    {
        "num": 2, "title": "AI Storyteller: EngSys Chronicles", "phase": "SPARK", "color": SPARK_COL,
        "date": "May 13, 2025", "type": "In-session (Groups)", "dur": "90 min",
        "desc": "Teams of 4 create illustrated AI stories from fictional EngSys scenarios — comic-strip style.",
        "seminar": "History & Evolution of AI",
        "presenter": "TBD (Random Draw)",
        "awards": ["Best Storyline", "Most Visual", "Best AI Use", "People's Choice"],
        "skills": ["Multi-turn prompting", "Visual storytelling", "Collaborative AI use"],
    },
    {
        "num": 3, "title": "Prompt Battle Arena", "phase": "SPARK", "color": SPARK_COL,
        "date": "Jun 10, 2025", "type": "In-session (Tournament)", "dur": "100 min",
        "desc": "Live prompt engineering tournament — Speed Round, Refinement Round, Domain Challenge, Bonus Round.",
        "seminar": "Prompt Engineering Deep Dive",
        "presenter": "TBD (Random Draw)",
        "awards": ["Speed King/Queen", "Best Improver", "Domain Master", "Overall Champion"],
        "skills": ["Personas & system prompts", "Few-shot / CoT", "Output formatting"],
    },
    {
        "num": 4, "title": "Copilot Power User", "phase": "BUILD", "color": BUILD_COL,
        "date": "Jul 8, 2025", "type": "Pre-worked + Demo", "dur": "100 min",
        "desc": "Each person automates 1 real recurring work task using M365 Copilot. 3-minute live demos.",
        "seminar": "M365 Copilot Ecosystem",
        "presenter": "TBD (Random Draw)",
        "awards": ["Most Time Saved", "Most Creative", "Best Before/After", "People's Choice"],
        "skills": ["M365 Copilot workflows", "Measuring AI ROI", "Daily tool integration"],
    },
    {
        "num": 5, "title": "Code with AI", "phase": "BUILD", "color": BUILD_COL,
        "date": "Aug 12, 2025", "type": "Pre-worked", "dur": "100 min",
        "desc": "Domain-specific coding tasks using AI. Annotate correct vs hallucinated code. Share working versions.",
        "seminar": "AI-Assisted Coding Best Practices",
        "presenter": "TBD (Random Draw)",
        "awards": ["Most Accurate Code", "Best Prompt", "Best Bug Catch", "Most Relevant"],
        "skills": ["AI code generation", "Code review discipline", "Hallucination awareness"],
    },
    {
        "num": 6, "title": "Smart Documentation / RAG", "phase": "BUILD", "color": BUILD_COL,
        "date": "Sep 9, 2025", "type": "Collaborative", "dur": "100 min",
        "desc": "Build a prototype AI Q&A system over team documentation. Test with domain-specific questions.",
        "seminar": "RAG — Retrieval-Augmented Generation",
        "presenter": "TBD (Random Draw)",
        "awards": ["Best Docs", "Best Questions", "Best Suggestion", "RAG Champion"],
        "skills": ["RAG fundamentals", "Document prep for AI", "Enterprise AI assistants"],
    },
    {
        "num": 7, "title": "Automate the Boring Stuff", "phase": "APPLY", "color": APPLY_COL,
        "date": "Oct 14, 2025", "type": "Pre-worked + Demo", "dur": "110 min",
        "desc": "Identify & prototype an AI automation for 1 repetitive daily task. 5-minute presentations.",
        "seminar": "Agentic AI — Agents vs Chatbots",
        "presenter": "TBD (Random Draw)",
        "awards": ["Most Impactful", "Most Creative", "Best Prototype", "Most Time Saved"],
        "skills": ["Automation opportunity ID", "AI solution design", "ROI estimation"],
    },
    {
        "num": 8, "title": "AI Agents Hackathon", "phase": "APPLY", "color": APPLY_COL,
        "date": "Nov 11, 2025", "type": "In-session (Groups)", "dur": "120 min",
        "desc": "Team hackathon with complex multi-step engineering scenarios. Design and implement agent workflows.",
        "seminar": "Building AI Agents — Frameworks & Patterns",
        "presenter": "TBD (Random Draw)",
        "awards": ["Best Architecture", "Most Realistic", "Best Collaboration", "Most Relevant"],
        "skills": ["Agent design patterns", "Multi-step reasoning", "Tool orchestration"],
    },
    {
        "num": 9, "title": "Data Detective", "phase": "APPLY", "color": APPLY_COL,
        "date": "Dec 9, 2025", "type": "Pre-worked + Demo", "dur": "100 min",
        "desc": "Analyse domain-specific datasets using AI. 5 NL questions, 1 visualization, 1 insight/anomaly.",
        "seminar": "AI for Data Analysis",
        "presenter": "TBD (Random Draw)",
        "awards": ["Best Insight", "Best Visualization", "Most Relevant", "Data Detective"],
        "skills": ["NL querying", "AI data visualization", "Anomaly detection"],
    },
    {
        "num": 10, "title": "Use Case: Problem Definition", "phase": "DELIVER", "color": DELIVER_COL,
        "date": "Jan 13, 2026", "type": "Pre-worked + Workshop", "dur": "110 min",
        "desc": "Capstone kickoff — each person identifies 1 real business use case. Formal 1-page proposals. 3-min pitches.",
        "seminar": "AI in Industrial Engineering — Case Studies",
        "presenter": "TBD (Random Draw)",
        "awards": ["Best Problem Def.", "Most Innovative", "Most Impactful", "Best Feasibility"],
        "skills": ["Business problem framing", "AI solution design", "Feasibility analysis"],
    },
    {
        "num": 11, "title": "Use Case: Build & Iterate", "phase": "DELIVER", "color": DELIVER_COL,
        "date": "Feb 10, 2026", "type": "Peer Review", "dur": "120 min",
        "desc": "Progress review — working prototypes, round-robin peer reviews, office hours, iteration planning.",
        "seminar": "Responsible AI — Ethics & Compliance",
        "presenter": "TBD (Random Draw)",
        "awards": ["Most Progress", "Best Prototype", "Best Feedback", "Most Resilient"],
        "skills": ["Iterative development", "Peer review", "Responsible AI"],
    },
    {
        "num": 12, "title": "Grand Finale Showcase", "phase": "DELIVER", "color": DELIVER_COL,
        "date": "Mar 10, 2026", "type": "Presentations", "dur": "150 min",
        "desc": "Final use case presentations to team & leadership. Judging panel. Annual leaderboard finale. Prizes!",
        "seminar": "Our AI Journey: 12 Months of Growth (Retrospective)",
        "presenter": "Vigneshvar SA",
        "awards": ["Best Overall (15pts)", "Runner-Up (10pts)", "2nd Runner-Up (7pts)", "People's Choice"],
        "skills": ["E2E value demonstration", "Presentation skills", "Stakeholder communication"],
    },
]

for md in months_detail:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Top bar with phase color
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.5), DARK)
    add_shape(slide, Inches(0), Inches(1.4), Inches(13.333), Inches(0.1), md['color'])

    # Month number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.3), Inches(1.0), Inches(0.8))
    badge.fill.solid()
    badge.fill.fore_color.rgb = md['color']
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = f"M{md['num']}"
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = 'Calibri'

    add_text_box(slide, Inches(2.0), Inches(0.3), Inches(9), Inches(0.5),
                 md['title'], font_size=30, color=WHITE, bold=True)

    info_text = f"{md['phase']}   ·   {md['date']}   ·   {md['dur']}   ·   {md['type']}"
    add_text_box(slide, Inches(2.0), Inches(0.85), Inches(9), Inches(0.4),
                 info_text, font_size=13, color=GREY_LIGHT)

    # Description card
    card = add_shape(slide, Inches(0.8), Inches(1.9), Inches(7.5), Inches(1.3), LIGHT)
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(7), Inches(0.3),
                 "ACTIVITY", font_size=11, color=md['color'], bold=True)
    add_text_box(slide, Inches(1.1), Inches(2.3), Inches(7), Inches(0.7),
                 md['desc'], font_size=14, color=DARK)

    # Seminar card
    card = add_shape(slide, Inches(0.8), Inches(3.4), Inches(7.5), Inches(1.1), LIGHT)
    add_text_box(slide, Inches(1.1), Inches(3.5), Inches(7), Inches(0.3),
                 "🎤 SEMINAR", font_size=11, color=PETROL, bold=True)
    add_text_box(slide, Inches(1.1), Inches(3.8), Inches(5), Inches(0.3),
                 md['seminar'], font_size=14, color=DARK, bold=True)
    add_text_box(slide, Inches(1.1), Inches(4.1), Inches(5), Inches(0.25),
                 f"Presenter: {md['presenter']}", font_size=12, color=GREY)

    # Awards
    card = add_shape(slide, Inches(9.0), Inches(1.9), Inches(3.8), Inches(2.6), LIGHT)
    add_text_box(slide, Inches(9.3), Inches(2.0), Inches(3.2), Inches(0.3),
                 "⭐ AWARDS", font_size=11, color=SPARK_COL, bold=True)
    add_multiline_box(slide, Inches(9.3), Inches(2.35), Inches(3.2), Inches(2.0),
                       md['awards'], font_size=13, color=DARK, bullet=True)

    # Skills
    card = add_shape(slide, Inches(0.8), Inches(4.8), Inches(12.0), Inches(1.2), LIGHT)
    add_text_box(slide, Inches(1.1), Inches(4.9), Inches(3), Inches(0.3),
                 "SKILLS DEVELOPED", font_size=11, color=PETROL, bold=True)
    for si, skill in enumerate(md['skills']):
        x = 1.1 + si * 3.8
        skill_badge = add_shape(slide, Inches(x), Inches(5.25), Inches(3.5), Inches(0.45), RGBColor(0xE0, 0xF7, 0xF5))
        add_text_box(slide, Inches(x + 0.2), Inches(5.3), Inches(3.1), Inches(0.35),
                     skill, font_size=12, color=PETROL_DARK)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 21 — THE DASHBOARD
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "The Dashboard — EMBRACE AI Web App", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.3), Inches(3))

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
             "A custom-built web application powering the entire initiative — accessible via a single URL, no login required for viewing.",
             font_size=15, color=GREY)

features_left = [
    ("Dashboard", "Live countdown, current activity, seminar spotlight, top 3 preview"),
    ("Leaderboard", "Monthly point breakdown, cumulative totals, avatar & domain tags"),
    ("Calendar", "12-month timeline with phase indicators and session details"),
    ("Team", "Member cards with photos, domains, roles, and domain filters"),
]
features_right = [
    ("Live Quiz", "Real-time MCQ with room codes, random order, tab-switch detection"),
    ("Survey", "Anonymous 1–5 star rating for seminar presenters"),
    ("Anti-Cheat", "Tab switch penalty (−1pt), copy/paste disabled during quiz"),
    ("Admin Panel", "Manage members, points, events, quizzes, surveys"),
]

for i, (title, desc) in enumerate(features_left):
    y = 2.3 + i * 1.05
    card = add_shape(slide, Inches(0.8), Inches(y), Inches(5.5), Inches(0.85), RGBColor(0x00, 0x20, 0x40))
    add_text_box(slide, Inches(1.1), Inches(y + 0.08), Inches(5), Inches(0.3),
                 title, font_size=15, color=GREEN, bold=True)
    add_text_box(slide, Inches(1.1), Inches(y + 0.4), Inches(5), Inches(0.35),
                 desc, font_size=12, color=GREY)

for i, (title, desc) in enumerate(features_right):
    y = 2.3 + i * 1.05
    card = add_shape(slide, Inches(6.8), Inches(y), Inches(5.5), Inches(0.85), RGBColor(0x00, 0x20, 0x40))
    add_text_box(slide, Inches(7.1), Inches(y + 0.08), Inches(5), Inches(0.3),
                 title, font_size=15, color=GREEN, bold=True)
    add_text_box(slide, Inches(7.1), Inches(y + 0.4), Inches(5), Inches(0.35),
                 desc, font_size=12, color=GREY)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             "Tech Stack: Python (Flask + SocketIO)  ·  Vanilla HTML/CSS/JS  ·  JSON Database  ·  Hosted at http://localhost:3000",
             font_size=12, color=GREY_LIGHT)

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 22 — EXPECTED OUTCOMES
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.3), DARK)
add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.7),
             "Expected Outcomes After 12 Months", font_size=36, color=WHITE, bold=True)

outcomes = [
    ("🎯", "GenAI-Ready Team", "All 16 members comfortable using SiemensGPT, M365 Copilot, and AI tools in daily work."),
    ("📦", "16 AI Use Cases", "Documented prototypes ready to pitch to US/Canada profit centers as value-adds."),
    ("⏱️", "Measurable Time Savings", "Aggregated productivity gains across the team — quantified ROI for management."),
    ("🧠", "AI-First Culture", "A mindset shift — team naturally considers AI solutions for new challenges."),
    ("📘", "Reusable Playbook", "The entire syllabus + dashboard can be shared with other Siemens teams in India."),
    ("🔗", "OT/AI Alignment", "Turns the team's formal AI engagement into demonstrated, documented capability."),
]

for i, (icon, title, desc) in enumerate(outcomes):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.8 + row * 1.7
    card = add_shape(slide, Inches(x), Inches(y), Inches(5.8), Inches(1.5), LIGHT)
    add_text_box(slide, Inches(x + 0.3), Inches(y + 0.15), Inches(0.6), Inches(0.5),
                 icon, font_size=24, color=DARK)
    add_text_box(slide, Inches(x + 0.9), Inches(y + 0.15), Inches(4.5), Inches(0.35),
                 title, font_size=17, color=DARK, bold=True)
    add_text_box(slide, Inches(x + 0.9), Inches(y + 0.55), Inches(4.5), Inches(0.7),
                 desc, font_size=13, color=GREY)

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 23 — 12-MONTH CALENDAR SUMMARY
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.3), DARK)
add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.7),
             "12-Month Calendar at a Glance", font_size=36, color=WHITE, bold=True)

# Header row
add_shape(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.5), PETROL)
headers = [("Mon", 0.5, 0.7), ("Date", 1.2, 1.5), ("Title", 2.7, 3.5), ("Phase", 6.2, 1.2),
           ("Type", 7.4, 2.0), ("Duration", 9.4, 1.2), ("Key Skill", 10.6, 2.2)]
for label, x, w in headers:
    add_text_box(slide, Inches(x), Inches(1.62), Inches(w), Inches(0.4),
                 label, font_size=11, color=WHITE, bold=True)

cal_rows = [
    (1, "Apr 8", "Who's Who: AI Edition", "SPARK", "Pre-worked", "90m", "Prompt basics"),
    (2, "May 13", "AI Storyteller", "SPARK", "In-session", "90m", "Multi-turn prompts"),
    (3, "Jun 10", "Prompt Battle Arena", "SPARK", "In-session", "100m", "Advanced prompting"),
    (4, "Jul 8", "Copilot Power User", "BUILD", "Pre + Demo", "100m", "M365 Copilot"),
    (5, "Aug 12", "Code with AI", "BUILD", "Pre-worked", "100m", "AI coding"),
    (6, "Sep 9", "Smart Documentation", "BUILD", "Collab", "100m", "RAG"),
    (7, "Oct 14", "Automate Boring Stuff", "APPLY", "Pre + Demo", "110m", "AI automation"),
    (8, "Nov 11", "AI Agents Hackathon", "APPLY", "In-session", "120m", "Agentic AI"),
    (9, "Dec 9", "Data Detective", "APPLY", "Pre + Demo", "100m", "Data analysis"),
    (10, "Jan 13", "Use Case: Define", "DELIVER", "Workshop", "110m", "Problem framing"),
    (11, "Feb 10", "Use Case: Build", "DELIVER", "Peer Rev", "120m", "Prototyping"),
    (12, "Mar 10", "Grand Finale", "DELIVER", "Present", "150m", "Value demo"),
]

phase_colors = {"SPARK": SPARK_COL, "BUILD": BUILD_COL, "APPLY": APPLY_COL, "DELIVER": DELIVER_COL}

for i, (mon, date, title, phase, typ, dur, skill) in enumerate(cal_rows):
    y = 2.15 + i * 0.42
    bg = LIGHT if i % 2 == 0 else WHITE
    add_shape(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(0.4), bg)
    vals = [(str(mon), 0.5, 0.7), (date, 1.2, 1.5), (title, 2.7, 3.5), (phase, 6.2, 1.2),
            (typ, 7.4, 2.0), (dur, 9.4, 1.2), (skill, 10.6, 2.2)]
    for vi, (val, x, w) in enumerate(vals):
        c = phase_colors.get(val, DARK) if vi == 3 else DARK
        b = vi == 3
        add_text_box(slide, Inches(x), Inches(y + 0.02), Inches(w), Inches(0.35),
                     val, font_size=11, color=c, bold=b)

slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 24 — THANK YOU
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text_box(slide, Inches(0.8), Inches(0.8), Inches(4), Inches(0.6),
             "SIEMENS", font_size=22, color=PETROL, bold=True)

add_text_box(slide, Inches(2.0), Inches(2.5), Inches(9), Inches(1.0),
             "Thank You", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.0), Inches(3.7), Inches(9), Inches(0.6),
             "Let's make this a fun and meaningful start to our AI learning journey.",
             font_size=18, color=GREY, alignment=PP_ALIGN.CENTER)

add_divider(slide, Inches(5.5), Inches(4.5), Inches(2.5))

add_text_box(slide, Inches(2.0), Inches(4.9), Inches(9), Inches(0.9),
             "Vigneshvar SA\nvigneshvar.sa@siemens.com\nEngineering Systems · Siemens Chennai",
             font_size=14, color=GREY_LIGHT, alignment=PP_ALIGN.CENTER, spacing=1.5)

# Decorative circles
for x, y, size, c in [(10.5, 5.0, 2.0, RGBColor(0x00,0x50,0x50)), (9.8, 1.0, 1.5, RGBColor(0x00,0x40,0x40))]:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()

add_text_box(slide, Inches(4.5), Inches(6.5), Inches(4.5), Inches(0.4),
             "vibe coded with AI ✦", font_size=11, color=PETROL, alignment=PP_ALIGN.CENTER)


# ═══ SAVE ═══
prs.save(OUT)
print(f"\n✅ Presentation saved to:\n   {OUT}\n")
print(f"   {len(prs.slides)} slides generated.")
